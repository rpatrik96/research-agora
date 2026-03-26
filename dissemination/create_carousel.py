"""Generate LinkedIn carousel PPTX for Research Agora."""

import io
import math

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colors ───────────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy  (#1a1a2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x4E, 0xCD, 0xC4)  # accent  (#4ecdc4)
CORAL = RGBColor(0xE9, 0x45, 0x60)  # secondary accent
MUTED = RGBColor(0xAA, 0xAA, 0xCC)  # muted white for body

# Hex strings for matplotlib (keep in sync with RGBColor constants above)
_DARK_BG_HEX = "#1a1a2e"
_TEAL_HEX = "#4ecdc4"

# ── Spiral brand colors (one per plugin category) ─────────────────────────────
# Order: outermost → innermost (matches generate_spiral_final.py SPIRAL_COLORS)
SPIRAL_COLORS = [
    "#4A90D9",  # academic       (outermost)
    "#50C878",  # development
    "#E8855A",  # editorial
    "#9B7FD4",  # formatting
    "#F5C242",  # office
    "#E05A7A",  # research-agents (innermost)
]

# ── Archimedean spiral geometry (mirrors generate_spiral_final.py) ────────────
_THETA_START = math.pi / 2
_TURNS = 2.0
_THETA_END = _THETA_START + _TURNS * 2 * math.pi
_R_MAX = 0.78

_raw_a = 0.04
_raw_b = 0.12
_raw_rmax = _raw_a + _raw_b * (_THETA_END - _THETA_START)
_scale = _R_MAX / _raw_rmax
_A = _raw_a * _scale
_B = _raw_b * _scale


def _archimedean_pts(n=2500):
    theta = np.linspace(_THETA_START, _THETA_END, n)
    r = _A + _B * (theta - _THETA_START)
    x = r * np.cos(theta)
    y = r * np.sin(theta)
    return np.column_stack([x, y])


def _arc_lengths(pts):
    d = np.diff(pts, axis=0)
    return np.concatenate([[0.0], np.cumsum(np.hypot(d[:, 0], d[:, 1]))])


def _color_segments(pts, colors):
    arcs = _arc_lengths(pts)
    total = arcs[-1]
    n = len(colors)
    segs = []
    for i, color in enumerate(colors):
        lo = total * i / n
        hi = total * (i + 1) / n
        mask = (arcs >= lo) & (arcs <= hi)
        idx = np.where(mask)[0]
        if len(idx) < 2:
            continue
        s = max(0, idx[0] - 1)
        e = min(len(pts) - 1, idx[-1] + 1)
        segs.append((pts[s : e + 1], color))
    return segs


def _draw_spiral(ax, lw=4.0):
    """Draw Archimedean spiral with category colors onto a matplotlib Axes."""
    pts = _archimedean_pts()
    segs = _color_segments(pts[::-1], SPIRAL_COLORS)
    for seg_pts, color in segs:
        ax.plot(
            seg_pts[:, 0],
            seg_pts[:, 1],
            color=color,
            linewidth=lw,
            solid_capstyle="round",
            solid_joinstyle="round",
            zorder=4,
        )
    # White center dot
    dot = plt.Circle((0, 0), 0.028, color="white", zorder=8)
    ax.add_patch(dot)


def render_spiral_png(size_px=400, dpi=100):
    """Render spiral mark to a PNG bytes buffer. Returns io.BytesIO."""
    size_in = size_px / dpi
    fig, ax = plt.subplots(figsize=(size_in, size_in), dpi=dpi)
    fig.patch.set_facecolor(_DARK_BG_HEX)
    ax.set_facecolor(_DARK_BG_HEX)
    half = 1.05
    ax.set_xlim(-half, half)
    ax.set_ylim(-half, half)
    ax.set_aspect("equal")
    ax.axis("off")
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    lw = max(2.0, size_in * dpi * 0.003)
    _draw_spiral(ax, lw=lw)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf

# ── Dimensions (1080×1080 px at 108 dpi → 10 × 10 inches) ───────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(10)

FONT_TITLE = 48
FONT_BODY = 26
FONT_SMALL = 20
FONT_TINY = 16

MARGIN_L = Inches(0.7)
CONTENT_W = Inches(8.6)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = BG_DARK):
    """Fill slide background with a solid color rectangle."""
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_accent_line(slide, y_pos=Inches(2.0), color=TEAL, width=Inches(1.5)):
    """Thin horizontal rule under title."""
    line = slide.shapes.add_shape(1, MARGIN_L, y_pos, width, Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_bottom_bar(slide, color=TEAL):
    """Colored bar at the very bottom."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(9.55), SLIDE_W, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_slide_number(slide, num: int, total: int = 9):
    txb = slide.shapes.add_textbox(Inches(8.5), Inches(9.6), Inches(1.3), Inches(0.35))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(FONT_TINY)
    run.font.color.rgb = BG_DARK
    run.font.name = "Calibri"
    run.font.bold = False


def add_title(slide, text: str, y=Inches(0.6), color=WHITE, size=FONT_TITLE):
    txb = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(1.1))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_body(slide, text: str, y=Inches(2.3), color=MUTED, size=FONT_BODY, height=Inches(6.5)):
    txb = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb


def _solid_shape(slide, x, y, w, h, color, border=False, border_color=None, border_pt=None):
    """Helper: add a filled rectangle, optionally with a border."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border and border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = border_pt or Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ============================================================================
# Slide builders
# ============================================================================


def slide1(prs):
    """Title slide — includes Archimedean spiral mark."""
    s = blank_slide(prs)
    fill_bg(s)

    # decorative top stripe
    _solid_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.35), TEAL)

    # ── Spiral mark (centered, upper half) ───────────────────────────────────
    # Render at 600×600 px for crisp display at 10-inch slide width.
    spiral_buf = render_spiral_png(size_px=600, dpi=150)
    spiral_size = Inches(2.8)
    spiral_x = (SLIDE_W - spiral_size) / 2  # horizontally centered
    spiral_y = Inches(0.55)
    s.shapes.add_picture(spiral_buf, spiral_x, spiral_y, spiral_size, spiral_size)

    # big title — centered, below spiral
    txb = s.shapes.add_textbox(MARGIN_L, Inches(3.55), CONTENT_W, Inches(1.4))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Research Agora"
    run.font.size = Pt(64)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # accent rule
    rule = s.shapes.add_shape(1, Inches(3.5), Inches(5.1), Inches(3.0), Pt(4))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    # subtitle
    txb2 = s.shapes.add_textbox(MARGIN_L, Inches(5.35), CONTENT_W, Inches(0.9))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "AI skills for researchers, by researchers."
    run2.font.size = Pt(28)
    run2.font.bold = False
    run2.font.color.rgb = TEAL
    run2.font.name = "Calibri"

    # tagline
    txb3 = s.shapes.add_textbox(MARGIN_L, Inches(6.45), CONTENT_W, Inches(0.7))
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Built by a researcher, for researchers."
    run3.font.size = Pt(FONT_SMALL)
    run3.font.bold = False
    run3.font.color.rgb = MUTED
    run3.font.name = "Calibri"
    run3.font.italic = True

    add_bottom_bar(s, CORAL)
    add_slide_number(s, 1)


def slide2(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), CORAL)

    add_title(s, "The Problem", y=Inches(0.6), color=CORAL, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=CORAL, width=Inches(2.5))

    body = (
        "Researchers use ChatGPT for everything.\n\n"
        "But there are no verified, reusable workflows.\n\n"
        "Citation hallucinations pass peer review.\n\n"
        "We need infrastructure — not just chat."
    )
    add_body(s, body, y=Inches(2.2), color=WHITE, size=FONT_BODY)

    add_bottom_bar(s, CORAL)
    add_slide_number(s, 2)


def slide3(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), TEAL)

    add_title(s, "Research Agora", y=Inches(0.6), color=TEAL, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=TEAL, width=Inches(3.0))

    body = (
        "A skills marketplace for ML research.\n\n"
        "61 public skills  \u00b7  6 plugins  \u00b7  MIT license\n\n"
        "One install command. Works today."
    )
    add_body(s, body, y=Inches(2.2), color=WHITE, size=FONT_BODY)

    add_bottom_bar(s, TEAL)
    add_slide_number(s, 3)


def slide4(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), TEAL)

    add_title(s, "What You Can Do", y=Inches(0.6), color=WHITE, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=TEAL, width=Inches(3.5))

    skills = [
        ("\u25b8 /paper-review", "Simulate a skeptical reviewer"),
        ("\u25b8 /paper-references", "Catch hallucinated citations"),
        ("\u25b8 /proof-auditor", "Verify mathematical proofs"),
        ("\u25b8 /claim-auditor", "Audit evidence for every claim"),
        ("\u25b8 /paper-introduction", "Draft compelling intros"),
    ]

    y = Inches(2.3)
    for cmd, desc in skills:
        cmd_box = s.shapes.add_textbox(MARGIN_L, y, Inches(3.8), Inches(0.55))
        tf = cmd_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = cmd
        run.font.size = Pt(FONT_BODY)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run.font.name = "Calibri"

        desc_box = s.shapes.add_textbox(Inches(4.6), y, Inches(4.8), Inches(0.55))
        tf2 = desc_box.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(FONT_BODY)
        run2.font.bold = False
        run2.font.color.rgb = MUTED
        run2.font.name = "Calibri"

        y += Inches(0.82)

    add_bottom_bar(s, TEAL)
    add_slide_number(s, 4)


def slide5(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), TEAL)

    add_title(s, "Verified, Not Vibes", y=Inches(0.6), color=WHITE, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=TEAL, width=Inches(4.0))

    txb = s.shapes.add_textbox(MARGIN_L, Inches(2.2), CONTENT_W, Inches(0.7))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Every verification skill ships with a benchmark."
    run.font.size = Pt(FONT_BODY)
    run.font.bold = False
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    levels = [
        ("L1", "Code-verified"),
        ("L2", "Reproducible experiment"),
        ("L3", "Paper evidence"),
        ("L4", "Citation support"),
        ("L5", "Logical argument"),
        ("L6", "Assertion (no evidence)"),
    ]
    y = Inches(3.2)
    for lvl, desc in levels:
        row = s.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.48))
        tf2 = row.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        r1 = p2.add_run()
        r1.text = f"{lvl}  "
        r1.font.size = Pt(22)
        r1.font.bold = True
        r1.font.color.rgb = TEAL
        r1.font.name = "Calibri"
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(22)
        r2.font.bold = False
        r2.font.color.rgb = MUTED
        r2.font.name = "Calibri"
        y += Inches(0.48)

    add_body(
        s,
        "\nYou can measure skill quality.",
        y=Inches(8.2),
        color=TEAL,
        size=FONT_SMALL,
        height=Inches(0.8),
    )

    add_bottom_bar(s, TEAL)
    add_slide_number(s, 5)


def slide6(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), CORAL)

    add_title(s, "Works Everywhere", y=Inches(0.6), color=WHITE, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=CORAL, width=Inches(3.5))

    platforms = [
        "Claude Code (native)",
        "Cursor (.mdc rules)",
        "Gemini CLI",
        "GitHub Copilot",
    ]
    y = Inches(2.5)
    for plat in platforms:
        row = s.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.7))
        tf = row.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = plat
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        y += Inches(0.85)

    _solid_shape(s, MARGIN_L, Inches(6.2), Inches(5.0), Pt(2), CORAL)

    add_body(
        s,
        "Conversion tool included.\nNo vendor lock-in.",
        y=Inches(6.5),
        color=MUTED,
        size=FONT_BODY,
        height=Inches(1.5),
    )

    add_bottom_bar(s, CORAL)
    add_slide_number(s, 6)


def slide7(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), TEAL)

    add_title(s, "Benchmarked", y=Inches(0.6), color=WHITE, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.75), color=TEAL, width=Inches(2.5))

    # big metric box
    metric_bg = s.shapes.add_shape(
        1, Inches(1.5), Inches(2.8), Inches(7.0), Inches(2.8)
    )
    metric_bg.fill.solid()
    metric_bg.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    metric_bg.line.color.rgb = TEAL
    metric_bg.line.width = Pt(1.5)

    txb_name = s.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(6.0), Inches(0.7))
    tf = txb_name.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "HALLMARK citation benchmark"
    run.font.size = Pt(22)
    run.font.bold = False
    run.font.color.rgb = MUTED
    run.font.name = "Calibri"

    txb_f1 = s.shapes.add_textbox(Inches(2.0), Inches(3.7), Inches(6.0), Inches(1.5))
    tf2 = txb_f1.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "F1-H: 0.908"
    run2.font.size = Pt(52)
    run2.font.bold = True
    run2.font.color.rgb = TEAL
    run2.font.name = "Calibri"

    add_body(
        s,
        "Quarterly refresh.\nOpen leaderboard.\n\nMeasure before you trust.",
        y=Inches(6.2),
        color=MUTED,
        size=FONT_BODY,
        height=Inches(2.5),
    )

    add_bottom_bar(s, TEAL)
    add_slide_number(s, 7)


def slide8(prs):
    s = blank_slide(prs)
    fill_bg(s)

    _solid_shape(s, Inches(0), Inches(0.0), Inches(0.18), Inches(2.1), CORAL)

    add_title(s, "P-AGI Workshop @ ICLR 2026", y=Inches(0.6), color=CORAL, size=36)
    add_accent_line(s, y_pos=Inches(1.75), color=CORAL, width=Inches(5.5))

    # paper title box
    paper_bg = s.shapes.add_shape(1, MARGIN_L, Inches(2.3), CONTENT_W, Inches(1.8))
    paper_bg.fill.solid()
    paper_bg.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    paper_bg.line.fill.background()

    txb = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(8.0), Inches(1.4))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "\u201cSkills, Benchmarks, and Verification\n"
        "Are What AI-Assisted Research Needs\u201d"
    )
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    run.font.italic = True

    body = (
        "Test-driven research:\n\n"
        "Define verification criteria\n"
        "BEFORE delegating to AI."
    )
    add_body(s, body, y=Inches(4.5), color=MUTED, size=FONT_BODY, height=Inches(3.0))

    add_bottom_bar(s, CORAL)
    add_slide_number(s, 8)


def slide9(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # full-width teal top stripe
    _solid_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.4), TEAL)

    add_title(s, "Get Started", y=Inches(0.7), color=WHITE, size=FONT_TITLE)
    add_accent_line(s, y_pos=Inches(1.85), color=TEAL, width=Inches(2.2))

    links = [
        ("Browse", "rpatrik96.github.io/research-agora"),
        ("Install", "/plugin marketplace add rpatrik96/research-agora"),
        ("Contribute", "github.com/rpatrik96/research-agora"),
    ]
    y = Inches(2.5)
    for label, val in links:
        lb = s.shapes.add_textbox(MARGIN_L, y, Inches(1.9), Inches(0.55))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label + ":"
        r.font.size = Pt(FONT_BODY)
        r.font.bold = True
        r.font.color.rgb = TEAL
        r.font.name = "Calibri"

        vb = s.shapes.add_textbox(Inches(2.7), y, Inches(7.0), Inches(0.55))
        tf2 = vb.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = val
        r2.font.size = Pt(FONT_BODY - 2)
        r2.font.bold = False
        r2.font.color.rgb = MUTED
        r2.font.name = "Calibri"

        y += Inches(1.0)

    _solid_shape(s, MARGIN_L, Inches(5.8), Inches(5.0), Pt(2), TEAL)

    # CTA bottom
    txb_cta = s.shapes.add_textbox(MARGIN_L, Inches(6.2), CONTENT_W, Inches(0.8))
    tf_cta = txb_cta.text_frame
    tf_cta.word_wrap = True
    p_cta = tf_cta.paragraphs[0]
    p_cta.alignment = PP_ALIGN.CENTER
    r_cta = p_cta.add_run()
    r_cta.text = "Save this post \u2193  \u00b7  Link in comments"
    r_cta.font.size = Pt(FONT_SMALL)
    r_cta.font.bold = True
    r_cta.font.color.rgb = WHITE
    r_cta.font.name = "Calibri"

    add_bottom_bar(s, TEAL)
    add_slide_number(s, 9)


# ============================================================================
# Main
# ============================================================================


def main():
    prs = new_prs()

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)

    out = "/Users/patrik.reizinger/Documents/GitHub/research-agora/dissemination/assets/linkedin-carousel.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
