#!/usr/bin/env python3
"""
Template Analyzer for Claude Skills

Extracts design specifications from PPTX files to create reusable style guides.
This allows Claude Code to follow consistent branding when creating new presentations.

Usage:
    python analyze_template.py path/to/template.pptx [--output slides|posters]

Output:
    Creates a STYLE.md file with extracted design specifications
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx required. Install with: pip install python-pptx")
    sys.exit(1)


def emu_to_inches(emu: int) -> float:
    """Convert EMUs to inches."""
    return emu / 914400


def emu_to_pt(emu: int) -> float:
    """Convert EMUs to points."""
    return emu / 12700


def rgb_to_hex(rgb: RGBColor) -> str:
    """Convert RGBColor to hex string."""
    if rgb is None:
        return None
    try:
        # Try different attribute patterns
        if hasattr(rgb, 'red'):
            return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
        elif hasattr(rgb, '__iter__'):
            r, g, b = rgb
            return f"#{r:02x}{g:02x}{b:02x}"
        elif isinstance(rgb, str):
            return rgb if rgb.startswith('#') else f"#{rgb}"
        else:
            # Try to convert from int representation
            val = int(rgb)
            return f"#{(val >> 16) & 0xFF:02x}{(val >> 8) & 0xFF:02x}{val & 0xFF:02x}"
    except Exception:
        return None


def extract_color_from_fill(fill) -> str | None:
    """Extract color from a fill object."""
    try:
        if fill.type is not None:
            if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                if fill.fore_color.rgb:
                    return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def analyze_presentation(pptx_path: str) -> dict[str, Any]:
    """Analyze a PowerPoint presentation and extract design specs."""
    prs = Presentation(pptx_path)

    specs = {
        "source_file": Path(pptx_path).name,
        "dimensions": {},
        "colors": {
            "backgrounds": [],
            "text": [],
            "accents": [],
        },
        "fonts": {
            "titles": [],
            "body": [],
            "all_used": set(),
        },
        "font_sizes": {
            "titles": [],
            "body": [],
            "captions": [],
        },
        "layouts": [],
        "slide_count": len(prs.slides),
    }

    # Slide dimensions
    specs["dimensions"] = {
        "width_inches": round(emu_to_inches(prs.slide_width), 2),
        "height_inches": round(emu_to_inches(prs.slide_height), 2),
        "aspect_ratio": "16:9" if abs(prs.slide_width / prs.slide_height - 16/9) < 0.1 else
                       "4:3" if abs(prs.slide_width / prs.slide_height - 4/3) < 0.1 else "custom",
    }

    # Analyze slide layouts
    for layout in prs.slide_layouts:
        layout_info = {
            "name": layout.name,
            "placeholders": []
        }
        for ph in layout.placeholders:
            layout_info["placeholders"].append({
                "type": str(ph.placeholder_format.type),
                "idx": ph.placeholder_format.idx,
            })
        specs["layouts"].append(layout_info)

    # Analyze each slide
    for slide_idx, slide in enumerate(prs.slides):
        # Background color
        if slide.background.fill.type is not None:
            bg_color = extract_color_from_fill(slide.background.fill)
            if bg_color and bg_color not in specs["colors"]["backgrounds"]:
                specs["colors"]["backgrounds"].append(bg_color)

        # Analyze shapes
        for shape in slide.shapes:
            # Text analysis
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Font name
                        if run.font.name:
                            specs["fonts"]["all_used"].add(run.font.name)

                        # Font size
                        if run.font.size:
                            size_pt = round(emu_to_pt(run.font.size), 0)

                            # Categorize by size
                            if size_pt >= 36:
                                if size_pt not in specs["font_sizes"]["titles"]:
                                    specs["font_sizes"]["titles"].append(size_pt)
                                if run.font.name and run.font.name not in specs["fonts"]["titles"]:
                                    specs["fonts"]["titles"].append(run.font.name)
                            elif size_pt >= 18:
                                if size_pt not in specs["font_sizes"]["body"]:
                                    specs["font_sizes"]["body"].append(size_pt)
                                if run.font.name and run.font.name not in specs["fonts"]["body"]:
                                    specs["fonts"]["body"].append(run.font.name)
                            else:
                                if size_pt not in specs["font_sizes"]["captions"]:
                                    specs["font_sizes"]["captions"].append(size_pt)

                        # Text color
                        try:
                            if run.font.color.rgb:
                                color = rgb_to_hex(run.font.color.rgb)
                                if color and color not in specs["colors"]["text"]:
                                    specs["colors"]["text"].append(color)
                        except Exception:
                            pass

            # Shape fill colors (accents)
            try:
                if hasattr(shape, 'fill') and shape.fill is not None:
                    color = extract_color_from_fill(shape.fill)
                    if color and color not in specs["colors"]["accents"]:
                        specs["colors"]["accents"].append(color)
            except Exception:
                pass

    # Convert set to list for JSON serialization
    specs["fonts"]["all_used"] = sorted(list(specs["fonts"]["all_used"]))

    # Sort font sizes
    specs["font_sizes"]["titles"] = sorted(specs["font_sizes"]["titles"], reverse=True)
    specs["font_sizes"]["body"] = sorted(specs["font_sizes"]["body"], reverse=True)
    specs["font_sizes"]["captions"] = sorted(specs["font_sizes"]["captions"], reverse=True)

    return specs


def generate_style_markdown(specs: dict, template_name: str) -> str:
    """Generate a STYLE.md file from extracted specifications."""
    # Pre-compute color values for the Python code block
    bg_color = (
        _hex_to_rgb_tuple(specs['colors']['backgrounds'][0])
        if specs['colors']['backgrounds']
        else '(255, 255, 255)'
    )
    text_color = (
        _hex_to_rgb_tuple(specs['colors']['text'][0])
        if specs['colors']['text']
        else '(0, 0, 0)'
    )
    accent_color = (
        _hex_to_rgb_tuple(specs['colors']['accents'][0])
        if specs['colors']['accents']
        else '(0, 102, 204)'
    )

    md = f"""# {template_name} Design Specifications

Auto-generated from `{specs['source_file']}` by analyze_template.py.

## Slide Dimensions

| Property | Value |
|----------|-------|
| Width | {specs['dimensions']['width_inches']}" |
| Height | {specs['dimensions']['height_inches']}" |
| Aspect Ratio | {specs['dimensions']['aspect_ratio']} |

## Color Palette

### Background Colors
{chr(10).join(f"- `{c}`" for c in specs['colors']['backgrounds']) or "- Not detected (likely white/transparent)"}

### Text Colors
{chr(10).join(f"- `{c}`" for c in specs['colors']['text'][:5]) or "- Default (black)"}

### Accent Colors
{chr(10).join(f"- `{c}`" for c in specs['colors']['accents'][:8]) or "- None detected"}

## Typography

### Fonts Used
| Purpose | Font(s) |
|---------|---------|
| Titles | {', '.join(specs['fonts']['titles'][:3]) or 'Default'} |
| Body | {', '.join(specs['fonts']['body'][:3]) or 'Default'} |
| All fonts | {', '.join(specs['fonts']['all_used'][:6]) or 'Default'} |

### Font Sizes
| Element | Sizes (pt) |
|---------|-----------|
| Titles | {', '.join(str(int(s)) for s in specs['font_sizes']['titles'][:4]) or 'N/A'} |
| Body | {', '.join(str(int(s)) for s in specs['font_sizes']['body'][:4]) or 'N/A'} |
| Captions | {', '.join(str(int(s)) for s in specs['font_sizes']['captions'][:4]) or 'N/A'} |

## Available Layouts

| Layout Name | Placeholders |
|-------------|--------------|
"""

    for layout in specs['layouts'][:10]:
        ph_types = [p['type'].replace('PLACEHOLDER_TYPE.', '') for p in layout['placeholders']]
        md += f"| {layout['name']} | {', '.join(ph_types[:4])} |\n"

    md += f"""

## Python-PPTX Configuration

```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Dimensions
SLIDE_WIDTH = Inches({specs['dimensions']['width_inches']})
SLIDE_HEIGHT = Inches({specs['dimensions']['height_inches']})

# Colors (as RGB tuples)
COLORS = {{
    "background": {bg_color},
    "text_primary": {text_color},
    "accent": {accent_color},
}}

# Font sizes
FONT_SIZES = {{
    "title": Pt({int(specs['font_sizes']['titles'][0]) if specs['font_sizes']['titles'] else 44}),
    "subtitle": Pt({int(specs['font_sizes']['titles'][1]) if len(specs['font_sizes']['titles']) > 1 else 32}),
    "body": Pt({int(specs['font_sizes']['body'][0]) if specs['font_sizes']['body'] else 24}),
    "caption": Pt({int(specs['font_sizes']['captions'][0]) if specs['font_sizes']['captions'] else 14}),
}}

# Fonts
FONTS = {{
    "title": "{specs['fonts']['titles'][0] if specs['fonts']['titles'] else 'Arial'}",
    "body": "{specs['fonts']['body'][0] if specs['fonts']['body'] else 'Arial'}",
}}
```

## Usage Notes

When creating presentations with this style:

1. **Start from template**: Copy the original .pptx file as your base
2. **Use slide layouts**: Prefer built-in layouts over manual positioning
3. **Match colors exactly**: Use the hex codes above for consistency
4. **Font hierarchy**: Maintain the size relationships for visual hierarchy
5. **Aspect ratio**: Ensure new slides match {specs['dimensions']['aspect_ratio']}

## Slide Count in Original

Total slides: {specs['slide_count']}
"""

    return md


def _hex_to_rgb_tuple(hex_color: str) -> str:
    """Convert hex color to RGB tuple string."""
    if not hex_color:
        return "(255, 255, 255)"
    hex_color = hex_color.lstrip('#')
    return f"({int(hex_color[0:2], 16)}, {int(hex_color[2:4], 16)}, {int(hex_color[4:6], 16)})"


def main():
    parser = argparse.ArgumentParser(
        description="Extract design specifications from PowerPoint templates"
    )
    parser.add_argument("pptx_file", help="Path to the .pptx template file")
    parser.add_argument(
        "--output", "-o",
        choices=["slides", "posters"],
        default="slides",
        help="Output directory (slides or posters)"
    )
    parser.add_argument(
        "--name", "-n",
        help="Template name (defaults to filename without extension)"
    )

    args = parser.parse_args()

    pptx_path = Path(args.pptx_file)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    if not pptx_path.suffix.lower() == '.pptx':
        print("Error: File must be a .pptx file")
        sys.exit(1)

    template_name = args.name or pptx_path.stem

    print(f"Analyzing: {pptx_path}")
    specs = analyze_presentation(str(pptx_path))

    # Determine output directory
    # Templates are stored in plugin directories:
    # - slides -> plugins/office/templates/slides/
    # - posters -> plugins/academic/templates/posters/
    script_dir = Path(__file__).parent.parent  # Go up from templates/ to repo root
    template_slug = template_name.lower().replace(" ", "-")
    if args.output == "slides":
        output_dir = script_dir / "plugins" / "office" / "templates" / "slides" / template_slug
    else:  # posters
        output_dir = script_dir / "plugins" / "academic" / "templates" / "posters" / template_slug
    output_dir.mkdir(parents=True, exist_ok=True)

    # Copy original template
    import shutil
    template_dest = output_dir / pptx_path.name
    shutil.copy2(pptx_path, template_dest)
    print(f"Copied template to: {template_dest}")

    # Write STYLE.md
    style_md = generate_style_markdown(specs, template_name)
    style_path = output_dir / "STYLE.md"
    style_path.write_text(style_md)
    print(f"Generated style guide: {style_path}")

    # Write raw specs as JSON for programmatic use
    specs_path = output_dir / "specs.json"
    # Convert sets to lists for JSON
    specs_json = json.dumps(specs, indent=2, default=str)
    specs_path.write_text(specs_json)
    print(f"Generated specs JSON: {specs_path}")

    print(f"\nDone! Template '{template_name}' added to templates/{args.output}/")
    print("Update paper-slides or paper-poster SKILL.md to reference this template.")


if __name__ == "__main__":
    main()
